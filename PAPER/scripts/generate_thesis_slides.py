# -*- coding: utf-8 -*-
"""
generate_thesis_slides.py
Completely redesigned for:
1. Concise, high-impact academic content (no walls of text, ~40 words per card)
2. Perfect vertical and horizontal centering (no overlapping, no off-center elements)
3. Elimination of the redundant/overlapping italic subtitle line on all slides
4. Seamless scenario support for SOICT 2026:
   - KLTN_slides_under_review.pptx (Default: Submitted & Under Review)
   - KLTN_slides_accepted_template.pptx (Prepared layout slot for Accepted scenario)
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Colors
C_PRIMARY = RGBColor(31, 78, 121)       # #1F4E79 Dark Navy
C_ACCENT = RGBColor(46, 117, 182)       # #2E75B6 Mid Blue
C_DARK = RGBColor(40, 40, 40)           # #282828 Charcoal body text
C_MUTED = RGBColor(110, 115, 125)       # #6E737D Muted gray
C_BG_CARD = RGBColor(246, 249, 252)     # #F6F9FC Soft light slate
C_BORDER_CARD = RGBColor(218, 226, 236) # #DAE2EC Subtle border
C_WHITE = RGBColor(255, 255, 255)
C_HIGHLIGHT_BG = RGBColor(254, 250, 235)# #FEFAEB Warm amber
C_HIGHLIGHT_BORDER = RGBColor(230, 180, 20)
C_SUCCESS_BG = RGBColor(235, 247, 238)  # #EBF7EE
C_SUCCESS_TEXT = RGBColor(27, 94, 32)
C_ALERT_BG = RGBColor(255, 238, 240)    # #FFEEF0
C_ALERT_TEXT = RGBColor(183, 28, 28)

FONT_FAMILY = "Arial"

# Images
IMG_PIPELINE = r'd:\KhoaLuanTotNghiep\PAPER\figures\Fig1_quanvolution_pipeline.png'
IMG_FEATURES = r'd:\KhoaLuanTotNghiep\PAPER\figures\Fig2_feature_comparison.png'
IMG_ARCH = r'd:\KhoaLuanTotNghiep\GD4\fig_software_architecture.png'
IMG_BREAST_BENCH = r'd:\KhoaLuanTotNghiep\PAPER\figures\Fig3_breastmnist_benchmark.png'
IMG_OCT_BENCH = r'd:\KhoaLuanTotNghiep\PAPER\figures\Fig3_octmnist_benchmark.png'
IMG_THETA = r'd:\KhoaLuanTotNghiep\PAPER\figures\Fig4c_theta_trajectories.png'
IMG_GRAD = r'd:\KhoaLuanTotNghiep\PAPER\figures\Fig4d_gradient_norms.png'
IMG_ABLATION = r'd:\KhoaLuanTotNghiep\GD2\results\figures\circuit_ablation_breastmnist.png'

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, tracker, action_title):
    """
    Standardized header:
    - Tracker label: y=0.38", h=0.25"
    - Action title: y=0.65", h=0.60" (Single concise takeaway)
    - Clean divider rule: y=1.30", h=0.015"
    - Total content area begins cleanly at y=1.45" to y=6.80" (Height: 5.35")
    """
    # Tracker
    tb_tr = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(11.733), Inches(0.25))
    tf_tr = tb_tr.text_frame
    tf_tr.word_wrap = True
    tf_tr.margin_left = tf_tr.margin_right = tf_tr.margin_top = tf_tr.margin_bottom = 0
    p_tr = tf_tr.paragraphs[0]
    p_tr.text = tracker.upper()
    p_tr.font.name = FONT_FAMILY
    p_tr.font.size = Pt(10)
    p_tr.font.bold = True
    p_tr.font.color.rgb = C_ACCENT

    # Action Title (Concise, single/double line, no overlapping)
    tb_t = slide.shapes.add_textbox(Inches(0.8), Inches(0.64), Inches(11.733), Inches(0.60))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
    p_t = tf_t.paragraphs[0]
    p_t.text = action_title
    p_t.font.name = FONT_FAMILY
    p_t.font.size = Pt(20)
    p_t.font.bold = True
    p_t.font.color.rgb = C_PRIMARY

    # Divider Rule
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.30), Inches(11.733), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = C_BORDER_CARD
    line.line.color.rgb = C_BORDER_CARD

def add_footer(slide, current_idx, total_slides=22, note="Trường ĐH Công nghệ Thông tin - ĐHQG-HCM"):
    # Divider Rule
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.92), Inches(11.733), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = C_BORDER_CARD
    line.line.color.rgb = C_BORDER_CARD

    # Note
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.98), Inches(9.0), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = f"Khóa luận Tốt nghiệp Kỹ sư Kỹ thuật Máy tính | {note}"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(9.5)
    p.font.color.rgb = C_MUTED

    # Page number
    tb_num = slide.shapes.add_textbox(Inches(10.5), Inches(6.98), Inches(2.033), Inches(0.3))
    tf_num = tb_num.text_frame
    tf_num.margin_left = tf_num.margin_right = tf_num.margin_top = tf_num.margin_bottom = 0
    p_num = tf_num.paragraphs[0]
    p_num.alignment = PP_ALIGN.RIGHT
    p_num.text = f"{current_idx} / {total_slides}"
    p_num.font.name = FONT_FAMILY
    p_num.font.size = Pt(10)
    p_num.font.bold = True
    p_num.font.color.rgb = C_PRIMARY

def add_clean_card(slide, left, top, width, height, title, bullet_items, bg_color=C_BG_CARD, border_color=C_BORDER_CARD, title_color=C_PRIMARY):
    """
    Creates a clean, perfectly centered card with concise bullet items.
    bullet_items: list of (bold_prefix, text) or single string.
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.2)

    tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.22), width - Inches(0.5), height - Inches(0.44))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    if title:
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = FONT_FAMILY
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = title_color
        p_t.space_after = Pt(10)

    for idx, item in enumerate(bullet_items):
        p = tf.add_paragraph() if (title or idx > 0) else tf.paragraphs[0]
        p.space_after = Pt(7)
        p.line_spacing = 1.15
        if isinstance(item, tuple):
            lead, body = item
            r_lead = p.add_run()
            r_lead.text = "• " + lead + (" " if not lead.endswith(":") else " ")
            r_lead.font.name = FONT_FAMILY
            r_lead.font.bold = True
            r_lead.font.size = Pt(11.5)
            r_lead.font.color.rgb = C_DARK

            r_body = p.add_run()
            r_body.text = body
            r_body.font.name = FONT_FAMILY
            r_body.font.size = Pt(11.5)
            r_body.font.color.rgb = C_DARK
        else:
            r = p.add_run()
            r.text = "• " + item
            r.font.name = FONT_FAMILY
            r.font.size = Pt(11.5)
            r.font.color.rgb = C_DARK

def add_stat_card(slide, left, top, width, height, big_num, label, desc="", num_color=C_PRIMARY, bg_color=C_BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = C_BORDER_CARD
    shape.line.width = Pt(1)

    tb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.12), width - Inches(0.2), height - Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    p_num = tf.paragraphs[0]
    p_num.alignment = PP_ALIGN.CENTER
    p_num.text = big_num
    p_num.font.name = FONT_FAMILY
    p_num.font.size = Pt(24)
    p_num.font.bold = True
    p_num.font.color.rgb = num_color

    p_lbl = tf.add_paragraph()
    p_lbl.alignment = PP_ALIGN.CENTER
    p_lbl.text = label
    p_lbl.font.name = FONT_FAMILY
    p_lbl.font.size = Pt(10.5)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = C_DARK
    p_lbl.space_before = Pt(2)

    if desc:
        p_d = tf.add_paragraph()
        p_d.alignment = PP_ALIGN.CENTER
        p_d.text = desc
        p_d.font.name = FONT_FAMILY
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = C_MUTED
        p_d.space_before = Pt(1)

def build_presentation(is_accepted=False, output_path=None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: TITLE SLIDE (Balanced Hero Layout)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1, C_PRIMARY)

    # University Header
    tb_u = s1.shapes.add_textbox(Inches(1.0), Inches(0.85), Inches(11.333), Inches(0.4))
    tf_u = tb_u.text_frame
    p_u = tf_u.paragraphs[0]
    p_u.text = "ĐẠI HỌC QUỐC GIA TP. HỒ CHÍ MINH  ·  TRƯỜNG ĐẠI HỌC CÔNG NGHỆ THÔNG TIN  ·  KHOA KỸ THUẬT MÁY TÍNH"
    p_u.font.name = FONT_FAMILY
    p_u.font.size = Pt(12)
    p_u.font.bold = True
    p_u.font.color.rgb = RGBColor(190, 215, 245)

    # Status Badge
    badge_bg = C_SUCCESS_BG if is_accepted else RGBColor(40, 95, 145)
    badge_text = "★ BÀI BÁO ĐÃ ĐƯỢC CHẤP NHẬN TẠI SOICT 2026 (SPRINGER CCIS)" if is_accepted else "BÀI BÁO HỘI NGHỊ QUỐC TẾ: ĐÃ NỘP & ĐANG BÌNH DUYỆT TẠI SOICT 2026 (SPRINGER CCIS)"
    badge_color = C_SUCCESS_TEXT if is_accepted else RGBColor(220, 235, 255)

    b_shape = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.45), Inches(8.2), Inches(0.36))
    b_shape.fill.solid()
    b_shape.fill.fore_color.rgb = badge_bg
    b_shape.line.color.rgb = badge_color
    tb_bt = s1.shapes.add_textbox(Inches(1.1), Inches(1.48), Inches(8.0), Inches(0.3))
    p_bt = tb_bt.text_frame.paragraphs[0]
    p_bt.text = badge_text
    p_bt.font.name = FONT_FAMILY
    p_bt.font.size = Pt(10.5)
    p_bt.font.bold = True
    p_bt.font.color.rgb = badge_color

    # Title Box
    tb_t = s1.shapes.add_textbox(Inches(1.0), Inches(2.05), Inches(11.333), Inches(2.0))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p_t1 = tf_t.paragraphs[0]
    p_t1.text = "NGHIÊN CỨU VÀ ỨNG DỤNG LỚP TÍCH CHẬP LƯỢNG TỬ (QUANVOLUTION)\nTRONG PHÂN LOẠI ẢNH Y TẾ MEDMNIST"
    p_t1.font.name = FONT_FAMILY
    p_t1.font.size = Pt(28)
    p_t1.font.bold = True
    p_t1.font.color.rgb = C_WHITE

    p_t2 = tf_t.add_paragraph()
    p_t2.text = "Đánh giá Thực nghiệm Đối xứng 1:1 giữa Mạch Lượng tử Tĩnh và Tự học trên Nền tảng Đa Hạt giống Khả Tái lập"
    p_t2.font.name = FONT_FAMILY
    p_t2.font.size = Pt(15)
    p_t2.font.italic = True
    p_t2.font.color.rgb = RGBColor(210, 228, 250)
    p_t2.space_before = Pt(8)

    # Accent divider
    line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.35), Inches(3.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = C_ACCENT
    line.line.color.rgb = C_ACCENT

    # Metadata Block
    tb_m = s1.shapes.add_textbox(Inches(1.0), Inches(4.65), Inches(11.333), Inches(1.8))
    tf_m = tb_m.text_frame
    p_m1 = tf_m.paragraphs[0]
    p_m1.text = "Báo cáo Khóa luận Tốt nghiệp Kỹ sư Ngành Kỹ thuật Máy tính"
    p_m1.font.name = FONT_FAMILY
    p_m1.font.size = Pt(14)
    p_m1.font.bold = True
    p_m1.font.color.rgb = C_WHITE

    p_m2 = tf_m.add_paragraph()
    p_m2.text = "• Sinh viên thực hiện: Nguyễn Hạo Nam (NamIsStudyingCE)\n• Giảng viên hướng dẫn: TS. Nguyễn Duy Xuân Bách\n• Thời gian & Địa điểm: TP. Hồ Chí Minh — Năm 2026"
    p_m2.font.name = FONT_FAMILY
    p_m2.font.size = Pt(12.5)
    p_m2.font.color.rgb = RGBColor(220, 230, 245)
    p_m2.space_before = Pt(6)
    p_m2.line_spacing = 1.25

    # =========================================================================
    # SLIDE 2: CLINICAL CONTEXT & DATA IMBALANCE
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2, C_WHITE)
    add_header(s2, "Phần 1: Bối cảnh & Động lực nghiên cứu",
               "Dữ liệu y tế khan hiếm và lệch lớp tự nhiên khiến Accuracy trở thành thước đo sai lầm")
    add_footer(s2, 2)

    # Left Card: Modalities
    add_clean_card(s2, Inches(0.8), Inches(1.55), Inches(5.7), Inches(5.1),
                   "Đặc thù Phương thái Hình ảnh Y tế",
                   [
                       ("Siêu âm tuyến vú (BUS):", "An toàn, không nhiễm xạ, tối ưu cho phụ nữ mô tuyến vú đặc; song chịu nhiễu đốm âm học (speckle noise) và biên độ tương phản thấp."),
                       ("Cắt lớp quang học võng mạc (OCT):", "Tiêu chuẩn vàng vi thể đáy mắt (CNV, DME, Drusen); khối lượng chụp diện rộng tạo áp lực đọc ảnh khổng lồ cho bác sĩ."),
                       ("Vai trò của hệ thống CAD:", "Đóng vai trò 'second opinion' khách quan, giảm tải sai sót do mệt mỏi thị giác và nắm bắt tương quan mô bệnh học cục bộ.")
                   ])

    # Right Card: Class Imbalance
    add_clean_card(s2, Inches(6.8), Inches(1.55), Inches(5.733), Inches(5.1),
                   "Bẫy Mất cân bằng Lớp & Rủi ro Lâm sàng",
                   [
                       ("Lệch lớp tự nhiên:", "Ca ác tính luôn chiếm thiểu số. Ví dụ BreastMNIST có tỷ lệ 73% lành tính / 27% ác tính."),
                       ("Thất bại của Accuracy:", "Một bộ phân loại đoán 100% 'Lành tính' vẫn đạt Accuracy 73%, nhưng bỏ sót 100% ca ác tính (Recall = 0%)."),
                       ("Chi phí bất đối xứng (Asymmetric Cost):", "Âm tính giả (False Negative — bỏ sót ung thư) gây hậu quả sinh mạng nghiêm trọng hơn nhiều so với Dương tính giả."),
                       ("Giải pháp phương pháp luận:", "Bắt buộc đánh giá đa chiều với Balanced Accuracy, F1 Macro, MCC và đặc biệt là PR-AUC — thước đo vàng cho lớp thiểu số.")
                   ])

    # =========================================================================
    # SLIDE 3: QUANTUM MACHINE LEARNING & QUANVOLUTION
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3, C_WHITE)
    add_header(s3, "Phần 1: Bối cảnh & Động lực nghiên cứu",
               "Mạng Quanvolution đề xuất cơ chế trích xuất đặc trưng phi tuyến trong không gian Hilbert 2^N chiều")
    add_footer(s3, 3)

    add_clean_card(s3, Inches(0.8), Inches(1.55), Inches(5.7), Inches(5.1),
                   "Nguyên lý Mạng Tích chập Lượng tử (Quanvolution)",
                   [
                       ("Khởi xướng (Henderson et al. 2020):", "Thay thế bộ lọc tích chập tuyến tính cổ điển bằng mạch lượng tử biến phân (VQC) trượt cục bộ trên ảnh."),
                       ("Không gian Hilbert 16 chiều:", "Mỗi patch 2×2 (4 pixel) được mã hóa vào 4 qubit; trạng thái lượng tử mở rộng ra không gian Hilbert 2⁴ = 16 chiều."),
                       ("Phép đo kỳ vọng Pauli-Z:", "Đo kỳ vọng trên từng qubit để thu về 4 giá trị thực [-1, 1], tạo thành 4 kênh feature map."),
                       ("Khác biệt với QCNN (Cong 2019):", "QCNN lượng tử hóa toàn mạng; Quanvolution là kiến trúc lai — mạch lượng tử trích xuất đặc trưng, phân loại do head cổ điển đảm nhiệm.")
                   ])

    add_clean_card(s3, Inches(6.8), Inches(1.55), Inches(5.733), Inches(5.1),
                   "Thiên kiến Quy nạp Lượng tử (Quantum Inductive Bias)",
                   [
                       ("Bản chất quy nạp (Kübler et al. 2021):", "Mạch lượng tử sở hữu inductive bias riêng, đặc biệt hiệu quả trên các tương quan topo phức tạp."),
                       ("Mệnh đề 'Power of Data' (Huang et al. 2021):", "Lợi thế lượng tử phụ thuộc chặt vào bản chất dữ liệu; dữ liệu lớn có thể xóa nhòa ưu thế của mạch lượng tử nông."),
                       ("Triết lý đánh giá (Schuld & Killoran 2022):", "Mục tiêu QML hiện nay không phải thổi phồng 'Quantum Advantage', mà là hiểu rõ khi nào và vì sao thành phần lượng tử có lợi."),
                       ("Đóng góp của đề tài:", "Lần đầu tiên phân định ranh giới hiệu quả theo 2 thái cực dữ liệu (nhỏ/lệch lớp vs lớn/đa lớp).")
                   ])

    # =========================================================================
    # SLIDE 4: THREE LITERATURE GAPS & SOLUTIONS
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4, C_WHITE)
    add_header(s4, "Phần 1: Bối cảnh & Động lực nghiên cứu",
               "Đề tài giải quyết 3 khoảng trống lớn của literature thông qua khung benchmark đối xứng 1:1")
    add_footer(s4, 4)

    w_card3 = Inches(3.64)
    add_clean_card(s4, Inches(0.8), Inches(1.55), w_card3, Inches(5.1),
                   "L1: Baseline Bất đối xứng",
                   [
                       ("Thực trạng:", "Literature so sánh với các CNN cổ điển tùy tiện, chênh lệch hàng nghìn tham số, gây sai lệch kết luận."),
                       ("Giải pháp đề tài:", "Thiết kế baseline CNN đối xứng 1:1: cùng đầu ra 784 chiều, cùng classifier head, chênh lệch toàn mạng đúng 20 tham số Conv2D."),
                       ("Ý nghĩa:", "Mọi chênh lệch hiệu năng đều quy về bản chất phép biến đổi đặc trưng.")
                   ], border_color=C_ACCENT)

    add_clean_card(s4, Inches(4.84), Inches(1.55), w_card3, Inches(5.1),
                   "L2: Đánh giá Thiếu Thống kê",
                   [
                       ("Thực trạng:", "Phần lớn công trình QML chỉ chạy 1–3 runs ngẫu nhiên, thiếu độ lệch chuẩn, khoảng tin cậy và kiểm định."),
                       ("Giải pháp đề tài:", "10 hạt giống cố định (S = {0, 42, ..., 2222}), kiểm định kép Paired t-test + Wilcoxon, đo Cohen's d và CI 95%."),
                       ("Ý nghĩa:", "Loại trừ triệt để yếu tố may rủi do khởi tạo ngẫu nhiên.")
                   ], border_color=C_ACCENT)

    add_clean_card(s4, Inches(8.88), Inches(1.55), w_card3, Inches(5.1),
                   "L3: Mạch Tĩnh vs Tự học",
                   [
                       ("Thực trạng:", "Tranh luận giữa cố định góc quay (Fixed) và học vi phân (Trainable) chưa được lượng hóa trên cùng bài toán."),
                       ("Giải pháp đề tài:", "Ma trận 3 tầng (Classical CNN, Fixed Quanv, Trainable Quanv) đo đạc chi phí phần cứng CPU chính xác."),
                       ("Ý nghĩa:", "Xác định rõ sự đánh đổi giữa dung lượng biểu diễn và chi phí tính toán.")
                   ], border_color=C_ACCENT)

    # =========================================================================
    # SLIDE 5: 1:1 SYMMETRICAL PIPELINE ARCHITECTURE
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5, C_WHITE)
    add_header(s5, "Phần 2: Phương pháp nghiên cứu",
               "Kiến trúc pipeline đối xứng 1:1 cố định hoàn toàn Classifier Head giữa hai phía")
    add_footer(s5, 5)

    if os.path.exists(IMG_PIPELINE):
        s5.shapes.add_picture(IMG_PIPELINE, Inches(0.8), Inches(1.65), width=Inches(6.4))

    add_clean_card(s5, Inches(7.5), Inches(1.55), Inches(5.033), Inches(5.1),
                   "Phân rã Tham số & Đối xứng Thiết kế",
                   [
                       ("Classifier Head cố định 100%:", "BatchNorm2d(4) → ReLU → Flatten → Linear(784 → K). Cấu trúc giống hệt nhau giữa Classical CNN và Quanvolution."),
                       ("BreastMNIST (K=2):", "Head Linear = 1,570 params (tổng Head+BN = 1,578). Toàn mạng CNN = 1,598 params. Chênh lệch đúng 20 tham số Conv2D."),
                       ("OCTMNIST (K=4):", "Head Linear = 3,140 params (tổng Head+BN = 3,148). Toàn mạng CNN = 3,168 params. Chênh lệch đúng 20 tham số Conv2D."),
                       ("Mạch Quanv Tĩnh:", "0 tham số học được (kernel hoàn toàn đóng băng)."),
                       ("Mạch Quanv Tự học:", "8 tham số (basic_L2) hoặc 24 tham số (strongly_L2) học đồng thời với head qua Adam lr kép."),
                       ("Kết luận:", "Loại bỏ hoàn toàn thiên kiến dung lượng tham số (Capacity Bias).")
                   ])

    # =========================================================================
    # SLIDE 6: QUANTUM ENCODING & THREE ANSATZ FAMILIES
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6, C_WHITE)
    add_header(s6, "Phần 2: Phương pháp nghiên cứu",
               "Ba họ ansatz khảo sát biểu diễn từ mạch đơn giản có cấu trúc đến mạch vướng víu toàn cục")
    add_footer(s6, 6)

    # Formula Box
    add_clean_card(s6, Inches(0.8), Inches(1.55), Inches(11.733), Inches(1.15),
                   "Quy trình Mã hóa Góc & Phép đo Kỳ vọng Pauli-Z",
                   [
                       ("Mã hóa góc (Angle Encoding):", "Mỗi patch 2×2 x = (x0, x1, x2, x3) ∈ [0, 1] được nạp vào trạng thái |ψ(x)⟩ = ⊗_i RY(π·x_i)|0⟩."),
                       ("Biến đổi & Phép đo:", "Qua ansatz U(θ), thu tensor đặc trưng: Fi(u, v) = ⟨Φ(x, θ)| Z_i |Φ(x, θ)⟩ ∈ [-1, 1] với i ∈ {0, 1, 2, 3}.")
                   ], bg_color=C_HIGHLIGHT_BG, border_color=C_HIGHLIGHT_BORDER)

    # 3 Ansatz Cards
    w_ansatz = Inches(3.64)
    add_clean_card(s6, Inches(0.8), Inches(2.9), w_ansatz, Inches(3.75),
                   "1. BasicEntanglerLayers",
                   [
                       ("Cấu trúc mạch:", "Cổng quay RY(θ) + vòng CNOT lân cận (qubit 0-1, 1-2, 2-3, 3-0)."),
                       ("Số tham số:", "4 tham số / tầng (với L=2: 8 tham số học được)."),
                       ("Đặc tính:", "Đơn giản, bảo toàn tính cục bộ, tối ưu cho dữ liệu quy mô nhỏ.")
                   ])

    add_clean_card(s6, Inches(4.84), Inches(2.9), w_ansatz, Inches(3.75),
                   "2. StronglyEntanglingLayers",
                   [
                       ("Cấu trúc mạch:", "Ba cổng quay tổng quát R(α, β, γ) trên mỗi qubit + CNOT mở rộng tầm xa."),
                       ("Số tham số:", "12 tham số / tầng (với L=2: 24 tham số học được)."),
                       ("Đặc tính:", "Không gian Hilbert giàu vướng víu và siêu vị, độ biểu đạt cao nhất.")
                   ])

    add_clean_card(s6, Inches(8.88), Inches(2.9), w_ansatz, Inches(3.75),
                   "3. RandomLayers",
                   [
                       ("Cấu trúc mạch:", "Cổng quay Haar-random phân bố đều kết hợp CNOT ngẫu nhiên cố định."),
                       ("Số tham số:", "0 tham số học được (kernel đóng băng hoàn toàn)."),
                       ("Đặc tính:", "Phép chiếu ngẫu nhiên phi tuyến; quán quân khảo sát GĐ2 trên OCT.")
                   ])

    # =========================================================================
    # SLIDE 7: SIX EVALUATION METRICS
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7, C_WHITE)
    add_header(s7, "Phần 2: Phương pháp nghiên cứu",
               "Hệ thống 6 chỉ số đánh giá đa chiều bảo đảm độ tin cậy trên bài toán y tế mất cân bằng")
    add_footer(s7, 7)

    add_clean_card(s7, Inches(0.8), Inches(1.55), Inches(5.7), Inches(5.1),
                   "Nhóm Chỉ số Phân loại Ngưỡng (Threshold Metrics)",
                   [
                       ("Accuracy = (TP + TN) / Tổng:", "Thước đo tổng thể; bị chi phối bởi lớp đa số (lành tính), chỉ mang tính tham khảo sơ bộ."),
                       ("Balanced Accuracy = (Sens + Spec) / 2:", "Trung bình cộng độ nhận diện đúng trên từng lớp; phản ánh công bằng hiệu năng lâm sàng."),
                       ("F1-score (Macro) = (1/K) · ∑ F1_k:", "Trung bình điều hòa giữa Precision và Recall; quy ước macro phạt nặng mô hình bỏ sót lớp thiểu số."),
                       ("MCC = (TP·TN − FP·FN) / √[...]:", "Hệ số tương quan Matthews [-1, +1]; chỉ số khắt khe nhất, chỉ cao khi dự đoán đúng cả 4 góc ma trận nhầm lẫn.")
                   ])

    add_clean_card(s7, Inches(6.8), Inches(1.55), Inches(5.733), Inches(5.1),
                   "Nhóm Chỉ số Xếp hạng Lâm sàng (Ranking Metrics)",
                   [
                       ("ROC-AUC (One-vs-Rest Macro):", "Đo diện tích dưới đường cong TPR vs FPR trên mọi ngưỡng; xác suất xếp hạng mẫu dương cao hơn mẫu âm ngẫu nhiên."),
                       ("Quy ước Đa lớp:", "Áp dụng cơ chế One-vs-Rest trung bình vĩ mô macro cho bài toán OCTMNIST 4 lớp bệnh lý."),
                       ("PR-AUC (Precision-Recall AUC):", "CHỈ SỐ VÀNG TRONG Y TẾ khi ca ác tính hiếm gặp; không bị làm sai lệch bởi số lượng lớn ca âm tính thật (TN)."),
                       ("Kiểm định Thống kê Kép:", "Paired Student's t-test song hành cùng Wilcoxon signed-rank test (α=0.05), chuẩn hóa bằng kích thước hiệu ứng Cohen's d.")
                   ])

    # =========================================================================
    # SLIDE 8: SOFTWARE ARCHITECTURE & GRADIENT SANITY CHECK
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8, C_WHITE)
    add_header(s8, "Phần 2: Phương pháp nghiên cứu",
               "Kiến trúc phần mềm 4 tầng và kiểm chứng đạo hàm giải tích đạt độ chính xác vật lý < 4.1e-8")
    add_footer(s8, 8)

    if os.path.exists(IMG_ARCH):
        s8.shapes.add_picture(IMG_ARCH, Inches(0.8), Inches(1.65), width=Inches(6.4))

    add_clean_card(s8, Inches(7.5), Inches(1.55), Inches(5.033), Inches(5.1),
                   "Thiết kế Hệ thống & Kiểm chứng Vật lý",
                   [
                       ("Kiến trúc 4 tầng chuẩn mực:", "Dữ liệu (medmnist_loader) → Mô hình (circuits/trainable_quanv) → Thí nghiệm (10 seeds run_gd3.py) → Đầu ra/Kiểm định."),
                       ("Bảo toàn tính toàn vẹn:", "Tự động kiểm tra tensor [B, 4, 14, 14], khoảng [-1, 1], bẫy lỗi NaN/Inf trong vi phân; cố định seed_everything đồng bộ."),
                       ("Kiểm chứng Đạo hàm Vật lý:", "Trước khi chạy 110 runs, gradient vi phân PyTorch (statevector) được đối chứng trực tiếp với Parameter-Shift Rule:"),
                       ("Công thức Shift-Rule:", "∂F / ∂θ = [F(θ + π/2) − F(θ − π/2)] / 2."),
                       ("Sai lệch đo được:", "|Δ| < 4.1 × 10⁻⁸ — cấp độ sai số float32/64, khẳng định thuật toán lan truyền ngược phản ánh chính xác 100% động học vật lý.")
                   ])

    # =========================================================================
    # SLIDE 9: BREASTMNIST BENCHMARK: FIXED BASIC WINS
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9, C_WHITE)
    add_header(s9, "Phần 3: Kết quả thực nghiệm & Đánh giá",
               "Trên dữ liệu nhỏ lệch lớp BreastMNIST: Mạch lượng tử tĩnh vượt trội Classical CNN")
    add_footer(s9, 9)

    if os.path.exists(IMG_BREAST_BENCH):
        s9.shapes.add_picture(IMG_BREAST_BENCH, Inches(0.8), Inches(1.65), width=Inches(6.4))

    add_clean_card(s9, Inches(7.5), Inches(1.55), Inches(5.033), Inches(5.1),
                   "Kết quả 10-Seed trên BreastMNIST (Bảng 4.1)",
                   [
                       ("Fixed Basic L2 đạt ROC-AUC cao nhất:", "0.8521 ± 0.0095 [0.8453, 0.8589], vượt CNN (0.8336 ± 0.0259) có ý nghĩa thống kê (p_ttest = 0.0298, p_wilcoxon = 0.0254, Cohen's d = +0.815)."),
                       ("Fixed Strongly L2 đạt PR-AUC cao nhất:", "0.9182 ± 0.0071 [0.9131, 0.9232], vượt CNN (0.9041 ± 0.0100) với ý nghĩa thống kê rất cao (p_ttest = 0.0023, p_wilcoxon = 0.0059, d = +1.332)."),
                       ("Ổn định phương sai gấp ~2.7 lần:", "Độ lệch chuẩn ROC-AUC của Fixed Basic là 0.0095 so với 0.0259 của CNN (tỷ số 2.7263×), kháng biến động hạt giống xuất sắc."),
                       ("Trainable Strongly đạt Balanced Acc cao nhất:", "0.6945 ± 0.0451 (+0.0344 so với Fixed Strongly, d = +0.677), khác biệt với CNN không có ý nghĩa thống kê (p = 0.6701)."),
                       ("Kết luận:", "Thêm tham số tự học không tự động mang lại hiệu năng cao hơn.")
                   ])

    # =========================================================================
    # SLIDE 10: BREASTMNIST MECHANISM: STRUCTURAL REGULARIZER
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s10, C_WHITE)
    add_header(s10, "Phần 3: Kết quả thực nghiệm & Đánh giá",
               "Cơ chế giải thích: Mạch tĩnh đóng vai trò như một bộ điều hòa cấu trúc (Structural Regularizer)")
    add_footer(s10, 10)

    add_clean_card(s10, Inches(0.8), Inches(1.55), Inches(5.7), Inches(5.1),
                   "Giả thuyết Điều hòa Cấu trúc (Structural Regularizer)",
                   [
                       ("Bối cảnh 546 mẫu huấn luyện:", "Trên tập dữ liệu nhỏ, CNN cổ điển với kernel tự do dễ khớp vào nhiễu đốm âm học của ảnh siêu âm."),
                       ("Miền giá trị bị chặn [-1, 1]:", "Mạch tĩnh ánh xạ patch ảnh qua phép đo Pauli-Z bị chặn nghiêm ngặt trong đoạn [-1, 1]."),
                       ("Thu hẹp không gian tìm kiếm:", "Head phía sau chỉ có 1,570 tham số tuyến tính, tiếp nhận đặc trưng cố định nên không gian tìm kiếm trọng số bị thu hẹp."),
                       ("Kháng biến động khởi tạo:", "Mô hình ít nhạy cảm với thay đổi hạt giống seed, giải thích vì sao độ lệch chuẩn giảm tới 2.7 lần.")
                   ])

    add_clean_card(s10, Inches(6.8), Inches(1.55), Inches(5.733), Inches(5.1),
                   "Bảo tồn Xếp hạng Lâm sàng (PR-AUC Ranking)",
                   [
                       ("Ưu thế Mạch Strongly-Entangling:", "Chứa đầy đủ cổng quay 3 trục và CNOT mở rộng, tạo bản đồ đặc trưng có tính phân tách phi tuyến cao."),
                       ("Xếp hạng mẫu thiểu số chuẩn xác:", "Đặc trưng vướng víu giúp head gán xác suất phân tách rõ ràng giữa u lành và u ác tính ở các ngưỡng quyết định cao."),
                       ("Biểu hiện Quantum Inductive Bias:", "Cấu trúc topo của mạch lượng tử, chứ không phải số lượng tham số, là yếu tố mang lại lợi thế."),
                       ("Hiệu quả thực tiễn:", "Đạt PR-AUC 0.9182 mà không cần tốn bất kỳ chi phí huấn luyện nào cho kernel (0 tham số tự học).")
                   ])

    # =========================================================================
    # SLIDE 11: OCTMNIST BENCHMARK: CLASSICAL CNN WINS
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s11, C_WHITE)
    add_header(s11, "Phần 3: Kết quả thực nghiệm & Đánh giá",
               "Trên dữ liệu lớn đa lớp OCTMNIST: Classical CNN áp đảo toàn diện mọi mô hình lượng tử")
    add_footer(s11, 11)

    if os.path.exists(IMG_OCT_BENCH):
        s11.shapes.add_picture(IMG_OCT_BENCH, Inches(0.8), Inches(1.65), width=Inches(6.4))

    add_clean_card(s11, Inches(7.5), Inches(1.55), Inches(5.033), Inches(5.1),
                   "Kết quả 10-Seed trên OCTMNIST (Bảng 4.2 & 4.3)",
                   [
                       ("Classical CNN dẫn đầu áp đảo cả 6 metrics:", "ROC-AUC 0.7505 ± 0.0240, PR-AUC 0.4991 ± 0.0297, Balanced Acc 0.4433 ± 0.0135, F1 Macro 0.3206 ± 0.0175, MCC 0.3156 ± 0.0198."),
                       ("Khoảng cách thống kê khổng lồ:", "CNN vượt mô hình lượng tử tốt nhất (Trainable Strongly) với Δ = +0.0583 ROC-AUC (p ≈ 0.0001, Cohen's d = +2.108) và Balanced Acc (d = +1.874)."),
                       ("Nguyên nhân Expressibility Bottleneck:", "Bài toán 4 lớp bệnh đáy mắt đòi hỏi nắm bắt chi tiết vi thể. Mạch 4-qubit nén 4 pixel về 4 số thực qua 1 tầng biến phân bị mất mát độ phân giải."),
                       ("Head cổ điển nhỏ:", "Linear(784 → 4) chỉ có 3,140 tham số, không đủ dung lượng bù đắp mất mát đặc trưng."),
                       ("Minh chứng cho 'Power of Data':", "Khi dữ liệu đủ lớn (5,000 ảnh), mô hình cổ điển tối ưu sẽ vượt trội mô hình lượng tử nông.")
                   ])

    # =========================================================================
    # SLIDE 12: OCTMNIST INTERNAL COMPARISON: TRAINABLE VS CHAMPION
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s12, C_WHITE)
    add_header(s12, "Phần 3: Kết quả thực nghiệm & Đánh giá",
               "So sánh nội bộ OCTMNIST: Tự học vượt mạch tĩnh cùng họ nhưng chỉ hòa với Quán quân tĩnh")
    add_footer(s12, 12)

    add_clean_card(s12, Inches(0.8), Inches(1.55), Inches(5.7), Inches(5.1),
                   "Tự học vs Mạch tĩnh trong Cùng Họ (Strongly)",
                   [
                       ("Ưu thế vi phân lượng tử nội bộ:", "Khi giữ nguyên cấu trúc Strongly, việc cập nhật góc θ giúp mô hình thích nghi tốt hơn với dữ liệu đa lớp."),
                       ("So sánh số liệu (Bảng 4.3):", "Trainable Strongly đạt ROC-AUC 0.6922 ± 0.0199 so với 0.6690 ± 0.0055 của Fixed Strongly."),
                       ("Ý nghĩa thống kê:", "Mức cải thiện Δ = +0.0232 có ý nghĩa thống kê rõ rệt (p_wilcoxon = 0.0098, Cohen's d = +1.050)."),
                       ("Đánh đổi tính toán:", "Đổi lại là chi phí tính toán tăng hơn 100 lần trong quá trình huấn luyện do phải backprop qua mạch lượng tử.")
                   ])

    add_clean_card(s12, Inches(6.8), Inches(1.55), Inches(5.733), Inches(5.1),
                   "Trainable Strongly vs Quán quân Tĩnh random_L1",
                   [
                       ("Quán quân tĩnh từ khảo sát GĐ2:", "Cấu hình random_L1 (0 tham số học) được chọn từ giai đoạn ablation đưa vào ma trận đối chứng."),
                       ("Kết quả so sánh trực diện:", "Trainable Strongly (0.6922 ± 0.0199) so với Quán quân tĩnh random_L1 (0.6912 ± 0.0071)."),
                       ("Khác biệt không có ý nghĩa:", "Khoảng cách chỉ là Δ = +0.0010 (p_ttest = 0.8875, Cohen's d = +0.046 — mức ảnh hưởng không đáng kể)."),
                       ("Độ ổn định vượt trội:", "Phương sai của random_L1 nhỏ hơn Trainable Strongly gần 3 lần (0.0071 so với 0.0199)."),
                       ("Bài học thiết kế:", "Chọn cấu trúc mạch tĩnh tối ưu mang lại hiệu năng tương đương tự học nhưng tiết kiệm hàng trăm lần chi phí.")
                   ])

    # =========================================================================
    # SLIDE 13: OPTIMIZATION DYNAMICS & GRADIENT SANITY CHECK
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s13, C_WHITE)
    add_header(s13, "Phần 3: Kết quả thực nghiệm & Đánh giá",
               "Động học gradient ổn định 0.2–0.5 và quỹ đạo góc trơn tru loại trừ nguy cơ Barren Plateaus")
    add_footer(s13, 13)

    if os.path.exists(IMG_THETA):
        s13.shapes.add_picture(IMG_THETA, Inches(0.8), Inches(1.55), width=Inches(5.7))

    if os.path.exists(IMG_GRAD):
        s13.shapes.add_picture(IMG_GRAD, Inches(6.8), Inches(1.55), width=Inches(5.733))

    add_clean_card(s13, Inches(0.8), Inches(5.25), Inches(11.733), Inches(1.45),
                   "Kết luận Động học Tối ưu",
                   [
                       ("Quỹ đạo góc quay θ(t) (Trái):", "Theo dõi 24 góc lượng tử qua 20 epochs; góc di chuyển có trật tự trong 10 epochs đầu và ổn định tại cực tiểu địa phương, không bị kẹt cứng hay dao động hỗn loạn."),
                       ("Chuẩn Gradient L2 (Phải):", "Chuẩn L2 trung bình duy trì quanh 0.2–0.5 (đỉnh từng seed đạt ~1.3), hoàn toàn cách xa ngưỡng suy biến gradient 10⁻⁵ đến 10⁻⁷, bác bỏ giả thuyết Barren Plateaus [14].")
                   ], bg_color=C_SUCCESS_BG, border_color=C_SUCCESS_TEXT)

    # =========================================================================
    # SLIDE 14: CIRCUIT ABLATION & COMPUTATIONAL COST
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s14, C_WHITE)
    add_header(s14, "Phần 3: Kết quả thực nghiệm & Đánh giá",
               "Khảo sát Ablation 6 cấu hình và phân tích chi phí: Precompute giúp tăng tốc 100 lần")
    add_footer(s14, 14)

    if os.path.exists(IMG_ABLATION):
        s14.shapes.add_picture(IMG_ABLATION, Inches(0.8), Inches(1.65), width=Inches(5.6))

    add_clean_card(s14, Inches(6.7), Inches(1.55), Inches(5.833), Inches(5.1),
                   "Phân tích Chi phí Phần cứng (Bảng 4.5)",
                   [
                       ("Độ trễ suy luận CPU (Intel):", "Conv cổ điển: 0.31 ms / ảnh; Quanvolution: 220 ms / ảnh (chênh lệch ~700 lần do phải mô phỏng 196 statevector mỗi ảnh)."),
                       ("Nghẽn mô phỏng (Simulation Bottleneck):", "Khoảng 99.98% tổng thời gian huấn luyện tiêu tốn vào việc mô phỏng mạch lượng tử trên CPU."),
                       ("Giải pháp Precompute cho Mạch tĩnh:", "Bản đồ đặc trưng được tính toán trước 1 lần duy nhất và lưu cache. Thời gian huấn luyện 10 seeds head phân loại giảm xuống chỉ còn ~18 giây!"),
                       ("Kết luận từ Ablation GĐ2:", "Trên dữ liệu nhỏ (Breast), mạch đơn giản basic_L2 (0.8508) vượt trội strongly (0.8104–0.8322). Ngược lại trên dữ liệu lớn (OCT), random_L1 đa dạng chiếu đặc trưng dẫn đầu (0.6905).")
                   ])

    # =========================================================================
    # SLIDE 15: THREE SCIENTIFIC CONCLUSIONS
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s15, C_WHITE)
    add_header(s15, "Phần 4: Kết luận & Đóng góp",
               "Ba kết luận khoa học đúc kết ranh giới hiệu quả thực nghiệm của Mạng Quanvolution")
    add_footer(s15, 15)

    w_c15 = Inches(3.64)
    add_clean_card(s15, Inches(0.8), Inches(1.55), w_c15, Inches(5.1),
                   "1. Phụ thuộc Chế độ Dữ liệu",
                   [
                       ("Định đề 'Power of Data':", "Hiệu quả mô hình lượng tử phụ thuộc chặt vào kích thước và độ phức tạp phân lớp của dữ liệu."),
                       ("Thắng thế rõ rệt:", "Trên dữ liệu nhỏ, lệch lớp (BreastMNIST), mạch lượng tử tĩnh vượt trội Classical CNN cả về ROC-AUC, PR-AUC và độ ổn định."),
                       ("Chạm trần biểu diễn:", "Trên dữ liệu lớn đa lớp (OCTMNIST), Classical CNN áp đảo tuyệt đối (Δ = +0.0583, d = +2.108).")
                   ], border_color=C_PRIMARY)

    add_clean_card(s15, Inches(4.84), Inches(1.55), w_c15, Inches(5.1),
                   "2. Vai trò Bộ Điều hòa Cấu trúc",
                   [
                       ("Mạch tĩnh 0 tham số:", "Kernel lượng tử cố định không tạo thêm bậc tự do, hạn chế tình trạng head khớp vào nhiễu hạt ảnh y tế."),
                       ("Ổn định phương sai ~2.7 lần:", "Độ lệch chuẩn ROC-AUC giảm từ 0.0259 xuống 0.0095."),
                       ("Tối ưu hóa PR-AUC lâm sàng:", "Không gian vướng víu giữ vững thứ tự phân loại các ca bệnh ác tính hiếm gặp mà không tốn chi phí học.")
                   ], border_color=C_PRIMARY)

    add_clean_card(s15, Inches(8.88), Inches(1.55), w_c15, Inches(5.1),
                   "3. Giới hạn của Tự học Vi phân",
                   [
                       ("Cải thiện cục bộ:", "Trong cùng họ mạch strongly, tối ưu hóa góc quay cải thiện ROC-AUC từ 0.6690 lên 0.6922 trên OCTMNIST."),
                       ("Không vượt mạch tĩnh tối ưu:", "Trainable chỉ hòa với Quán quân tĩnh random_L1 (0.6912, p = 0.8875, d = +0.046)."),
                       ("Đánh đổi thực tiễn:", "Ablation chọn mạch tĩnh là giải pháp tối ưu chi phí tính toán cho ứng dụng thực tế.")
                   ], border_color=C_PRIMARY)

    # =========================================================================
    # SLIDE 16: FOUR CONTRIBUTIONS & SOICT PUBLICATION
    # =========================================================================
    s16 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s16, C_WHITE)
    add_header(s16, "Phần 4: Kết luận & Đóng góp",
               "Bốn đóng góp mới của đề tài và vị thế công bố khoa học quốc tế tại SOICT 2026")
    add_footer(s16, 16)

    add_clean_card(s16, Inches(0.8), Inches(1.55), Inches(5.7), Inches(5.1),
                   "Ba Đóng góp Phương pháp luận & Thực nghiệm",
                   [
                       ("C1. Khung Benchmark Đối xứng 1:1:", "Khắc phục triệt để 3 khoảng trống L1–L3 của literature; chuẩn hóa dung lượng tham số và kiểm định kép Paired t-test + Wilcoxon trên 10 seeds."),
                       ("C2. Bản đồ Phân định Data-Regime:", "Cung cấp bằng chứng thực nghiệm minh bạch đầu tiên về điểm mạnh (dữ liệu nhỏ lệch lớp) và điểm nghẽn biểu diễn (dữ liệu lớn đa lớp) của quanvolution 4-qubit."),
                       ("C3. Hệ thống Tái lập Hoàn toàn & Demo:", "Mã nguồn mở chuẩn hóa gắn git tag soict-submission-v4, đối chứng đạo hàm < 4.1e-8, và notebook demo chạy live sai lệch chỉ 1.47e-8.")
                   ])

    if is_accepted:
        add_clean_card(s16, Inches(6.8), Inches(1.55), Inches(5.733), Inches(5.1),
                       "★ C4. Bài báo Khoa học Quốc tế (Đã Chấp nhận Đăng)",
                       [
                           ("[KỊCH BẢN ĐÃ ACCEPT — KHUNG THIẾT KẾ SẴN ĐỂ ĐIỀN THÔNG TIN]", ""),
                           ("Tên bài báo:", "'Symmetrical Empirical Evaluation of Trainable versus Fixed Quanvolutional Filters in Medical Image Classification: A Rigorous, Reproducible Benchmark on MedMNIST'."),
                           ("Hội nghị công bố:", "The 13th International Symposium on Information and Communication Technology (SOICT 2026)."),
                           ("Nhà xuất bản:", "Springer CCIS (Communications in Computer and Information Science)."),
                           ("Vị thế học thuật:", "Được hội đồng phản biện quốc tế thẩm định độc lập và chấp nhận xuất bản chính thức."),
                           ("(Ghi chú chuẩn bị):", "Khung này đã sẵn sàng để điền: Số trang kỷ yếu, ngày thuyết trình tại hội nghị, và các đánh giá nổi bật của Reviewer ngay khi có thông báo chính thức.")
                       ], bg_color=C_SUCCESS_BG, border_color=C_SUCCESS_TEXT, title_color=C_SUCCESS_TEXT)
    else:
        add_clean_card(s16, Inches(6.8), Inches(1.55), Inches(5.733), Inches(5.1),
                       "C4. Bài báo Khoa học Quốc tế (Đã nộp & Đang bình duyệt)",
                       [
                           ("Tên bài báo khoa học:", "'Symmetrical Empirical Evaluation of Trainable versus Fixed Quanvolutional Filters in Medical Image Classification: A Rigorous, Reproducible Benchmark on MedMNIST'."),
                           ("Hội nghị nộp:", "Hội nghị Quốc tế SOICT 2026 (The 13th International Symposium on Information and Communication Technology)."),
                           ("Ấn phẩm xuất bản:", "Kỷ yếu Springer CCIS (Communications in Computer and Information Science)."),
                           ("Quy cách & Tính toàn vẹn:", "Bài báo dài 13 trang định dạng chuẩn Springer llncs; bản PDF nộp được đóng băng tại git tag soict-submission-v4 và mã băm SHA-256 xác thực c8d5d093..."),
                           ("Trạng thái hiện tại:", "Đã nộp thành công qua EasyChair và đang trong quá trình bình duyệt đồng đẳng (Peer Review).")
                       ], bg_color=C_HIGHLIGHT_BG, border_color=C_HIGHLIGHT_BORDER)

    # =========================================================================
    # SLIDE 17: FOUR FUTURE DIRECTIONS
    # =========================================================================
    s17 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s17, C_WHITE)
    add_header(s17, "Phần 4: Kết luận & Đóng góp",
               "Bốn hướng phát triển tiếp theo nhằm mở rộng quy mô và kiểm thử phần cứng vật lý")
    add_footer(s17, 17)

    w_c17 = Inches(2.7)
    add_clean_card(s17, Inches(0.8), Inches(1.55), w_c17, Inches(5.1),
                   "1. Tăng tốc GPU",
                   [
                       ("Mục tiêu:", "Vượt qua rào cản tốc độ mô phỏng CPU."),
                       ("Giải pháp:", "Tích hợp backend NVIDIA cuQuantum và PennyLane-Lightning-GPU."),
                       ("Mở rộng:", "Mở rộng kernel lên 3×3 (9-qubit) và 4×4 (16-qubit) kiểm tra khả năng vượt qua expressibility bottleneck trên OCT.")
                   ])

    add_clean_card(s17, Inches(3.8), Inches(1.55), w_c17, Inches(5.1),
                   "2. Kiểm thử NISQ Thực tế",
                   [
                       ("Mục tiêu:", "Đánh giá mô hình dưới tác động của nhiễu vật lý."),
                       ("Giải pháp:", "Triển khai mạch trên bộ xử lý thực IBM Quantum (Eagle/Heron)."),
                       ("Giảm thiểu lỗi (QEM):", "Ứng dụng Zero-Noise Extrapolation (ZNE) bảo toàn chất lượng feature map trước nhiễu depolarizing và readout.")
                   ])

    add_clean_card(s17, Inches(6.8), Inches(1.55), w_c17, Inches(5.1),
                   "3. Mở rộng MedMNIST",
                   [
                       ("Mục tiêu:", "Khảo sát phổ dữ liệu y sinh toàn diện."),
                       ("Giải pháp:", "Mở rộng benchmark sang toàn bộ 12 datasets 2D (Pneumonia, Path, Derma...)."),
                       ("Dữ liệu 3D:", "Khảo sát trên ảnh cắt lớp thể tích 3D (OrganMNIST3D, NoduleMNIST3D) với kernel lượng tử 3 chiều.")
                   ])

    add_clean_card(s17, Inches(9.8), Inches(1.55), w_c17, Inches(5.1),
                   "4. Kiến trúc Lai Tiên tiến",
                   [
                       ("Mục tiêu:", "Nâng cao tính biểu đạt và tương thích sâu."),
                       ("Giải pháp:", "Nghiên cứu cơ chế mã hóa Haar-random (HUAR) hoặc Hamiltonian encoding."),
                       ("Cơ chế thích ứng:", "Tích hợp Quantum Attention và Adaptive Quantum Gating để tự động điều tiết tỷ trọng đặc trưng lượng tử.")
                   ])

    # =========================================================================
    # SLIDE 18: DEFENSE SUMMARY & Q&A
    # =========================================================================
    s18 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s18, C_WHITE)
    add_header(s18, "Phần 5: Tổng kết & Bảo vệ",
               "Khóa luận hoàn thành toàn diện 4 mục tiêu nghiên cứu với cam kết liêm chính học thuật cao nhất")
    add_footer(s18, 18, note="Phiên thảo luận & Hỏi đáp (Q&A)")

    w_c18 = Inches(3.64)
    add_clean_card(s18, Inches(0.8), Inches(1.55), w_c18, Inches(3.2),
                   "Symmetrical Benchmark",
                   [
                       ("Thiết kế 1:1 chuẩn tắc:", "Cố định classifier head, chuẩn hóa dung lượng tham số chênh lệch đúng 20 params conv cổ điển."),
                       ("Đánh giá khắt khe:", "10 seeds độc lập, kiểm định thống kê kép t-test + Wilcoxon, CI 95% và Cohen's d.")
                   ], border_color=C_PRIMARY)

    add_clean_card(s18, Inches(4.84), Inches(1.55), w_c18, Inches(3.2),
                   "Empirical Boundary",
                   [
                       ("Thắng lợi rõ rệt:", "Mạch tĩnh vượt Classical CNN trên dữ liệu nhỏ lệch lớp BreastMNIST (ROC-AUC 0.8521, PR-AUC 0.9182, std giảm 2.7×)."),
                       ("Điểm nghẽn biểu diễn:", "CNN áp đảo trên dữ liệu lớn OCTMNIST (Δ = +0.0583, d = +2.108).")
                   ], border_color=C_PRIMARY)

    add_clean_card(s18, Inches(8.88), Inches(1.55), w_c18, Inches(3.2),
                   "Reproducible Artifacts",
                   [
                       ("Mã nguồn & Dữ liệu:", "100% tái lập được với seed cố định, tag soict-submission-v4."),
                       ("Bài báo Quốc tế:", "13 trang Springer CCIS nộp tại SOICT 2026."),
                       ("Hệ thống Demo:", "Live notebook kiểm thử tức thì với độ sai lệch 1.47e-8.")
                   ], border_color=C_PRIMARY)

    add_clean_card(s18, Inches(0.8), Inches(4.95), Inches(11.733), Inches(1.7),
                   "Thông tin Lưu trữ & Mã nguồn Dự án",
                   [
                       ("Kho lưu trữ GitHub chính thức:", "https://github.com/NamIsStudyingCE/Quanvolution (Gắn tag phiên bản: soict-submission-v4)"),
                       ("Bộ tài liệu nghiệm thu:", "Luận văn 58 trang (KLTN_draft_full.docx), Bảng số liệu Canonical (reconciliation_canonical.json), Video demo bảo vệ 104s."),
                       ("Lời cảm ơn:", "Chân thành cảm ơn TS. Nguyễn Duy Xuân Bách và Quý Thầy Cô Khoa Kỹ thuật Máy tính đã tận tình hướng dẫn và định hướng nghiên cứu.")
                   ], bg_color=C_HIGHLIGHT_BG, border_color=C_HIGHLIGHT_BORDER)

    # =========================================================================
    # SLIDE 19: BACKUP 1 - STATISTICAL SIGNIFICANCE MATRIX
    # =========================================================================
    s19 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s19, C_WHITE)
    add_header(s19, "Phụ lục Phòng vệ 1: Thống kê & Kiểm định Giả thuyết",
               "Ma trận kiểm định thống kê kép và kích thước hiệu ứng Cohen's d cho các cặp so sánh then chốt")
    add_footer(s19, 19, note="Slide Phụ trợ Trả lời Phản biện (Backup Slide 1)")

    add_clean_card(s19, Inches(0.8), Inches(1.55), Inches(11.733), Inches(5.1),
                   "Chi tiết Bảng Kiểm định Thống kê (Bảng 4.3 Canonical)",
                   [
                       ("1. Breast: Classical CNN vs Fixed Basic (AUC):", "Δ = −0.0186 | p_ttest = 0.0298 | p_wilcoxon = 0.0254 | Cohen's d = −0.815 (Fixed Basic thắng có ý nghĩa thống kê, hiệu ứng lớn)."),
                       ("2. Breast: Classical CNN vs Fixed Strongly (PR):", "Δ = −0.0140 | p_ttest = 0.0023 | p_wilcoxon = 0.0059 | Cohen's d = −1.332 (Fixed Strongly thắng áp đảo ở PR-AUC, hiệu ứng rất lớn)."),
                       ("3. Breast: Fixed Strongly vs Trainable Strongly (BAcc):", "Δ = −0.0344 | p_ttest = 0.0611 | p_wilcoxon = 0.0879 | Cohen's d = −0.677 (Mức xu hướng, chưa đạt p < 0.05)."),
                       ("4. OCT: Classical CNN vs Trainable Strongly (AUC):", "Δ = +0.0583 | p_ttest = 0.0001 | p_wilcoxon = 0.0020 | Cohen's d = +2.108 (CNN cổ điển áp đảo tuyệt đối, hiệu ứng khổng lồ)."),
                       ("5. OCT: Trainable Strongly vs Fixed Strongly (AUC):", "Δ = +0.0232 | p_ttest = 0.0090 | p_wilcoxon = 0.0098 | Cohen's d = +1.050 (Tự học cải thiện có ý nghĩa so với mạch tĩnh cùng họ)."),
                       ("6. OCT: Trainable Strongly vs Quán quân tĩnh random_L1 (AUC):", "Δ = +0.0010 | p_ttest = 0.8875 | p_wilcoxon = 0.6250 | Cohen's d = +0.046 (Khác biệt hoàn toàn không đáng kể, p = 0.8875)."),
                       ("Lưu ý về ngưỡng p-value Wilcoxon (n=10):", "Phân phối Wilcoxon signed-rank rời rạc với n=10 có mức tối thiểu lý thuyết là 1 / 2⁹ ≈ 0.00195. Mức p ≈ 0.0020 là mức ý nghĩa tối đa có thể đạt được.")
                   ])

    # =========================================================================
    # SLIDE 20: BACKUP 2 - PARAMETER SYMMETRY DECOMPOSITION
    # =========================================================================
    s20 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s20, C_WHITE)
    add_header(s20, "Phụ lục Phòng vệ 2: Phân rã Tham số Mô hình",
               "Chứng minh toán học về tính đối xứng 1:1 giữa Classical CNN và Quanvolution")
    add_footer(s20, 20, note="Slide Phụ trợ Trả lời Phản biện (Backup Slide 2)")

    add_clean_card(s20, Inches(0.8), Inches(1.55), Inches(5.7), Inches(5.1),
                   "Phân rã Tham số BreastMNIST (K=2)",
                   [
                       ("Classifier Head chung:", "Linear(784 → 2) có 784 × 2 weights + 2 biases = 1,570 tham số. BatchNorm2d(4) có 8 tham số. Tổng Head = 1,578 tham số."),
                       ("Mô hình Classical CNN:", "Lớp Conv2D(1→4, 2×2, bias) có (1×2×2 + 1) × 4 = 20 tham số. Tổng toàn mạng = 20 + 1,578 = 1,598 tham số."),
                       ("Mô hình Quanvolution Tĩnh:", "Kernel lượng tử = 0 tham số. Tổng toàn mạng = 1,578 tham số. Chênh lệch đúng 20 tham số của lớp Conv2D."),
                       ("Mô hình Quanvolution Tự học:", "Basic_L2 có 8 tham số góc (Tổng 1,586); Strongly_L2 có 24 tham số góc (Tổng 1,602 tham số)."),
                       ("Độ chênh lệch tối đa:", "Không vượt quá 24 tham số so với toàn mạng ~1,600 tham số (< 1.5%).")
                   ])

    add_clean_card(s20, Inches(6.8), Inches(1.55), Inches(5.733), Inches(5.1),
                   "Phân rã Tham số OCTMNIST (K=4)",
                   [
                       ("Classifier Head chung:", "Linear(784 → 4) có 784 × 4 weights + 4 biases = 3,140 tham số. BatchNorm2d(4) có 8 tham số. Tổng Head = 3,148 tham số."),
                       ("Mô hình Classical CNN:", "Lớp Conv2D(1→4, 2×2, bias) có 20 tham số. Tổng toàn mạng = 20 + 3,148 = 3,168 tham số."),
                       ("Mô hình Quanvolution Tĩnh:", "Kernel lượng tử = 0 tham số. Tổng toàn mạng = 3,148 tham số. Chênh lệch đúng 20 tham số."),
                       ("Mô hình Quanvolution Tự học:", "Basic_L1 có 4 tham số góc (Tổng 3,152); Strongly_L1 có 12 tham số góc (Tổng 3,160 tham số)."),
                       ("Quy chuẩn nghiêm ngặt:", "TUYỆT ĐỐI không có sự nhầm lẫn giữa 1,570 tham số (K=2) và 3,140 tham số (K=4).")
                   ])

    # =========================================================================
    # SLIDE 21: BACKUP 3 - PARAMETER-SHIFT RULE VS BACKPROP (Fixed User Overlap!)
    # =========================================================================
    s21 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s21, C_WHITE)
    # The action title is single, clean, concise: NO overlapping subtitle!
    add_header(s21, "Phụ lục Phòng vệ 3: Kiểm chứng Đạo hàm Lượng tử",
               "Chứng minh tính tương đương vật lý giữa Statevector Backpropagation và Parameter-Shift Rule")
    add_footer(s21, 21, note="Slide Phụ trợ Trả lời Phản biện (Backup Slide 3)")

    # 2 balanced columns instead of 1 dense block!
    add_clean_card(s21, Inches(0.8), Inches(1.55), Inches(5.7), Inches(5.1),
                   "Cơ sở Toán học & Hai Đường Tính Gradient",
                   [
                       ("1. Đạo hàm giải tích Statevector (PyTorch):", "Sử dụng engine lan truyền ngược PyTorch tích hợp PennyLane (Bergholm et al. 2018) trên không gian vector trạng thái mô phỏng."),
                       ("2. Parameter-Shift Rule (Schuld et al. 2019):", "Công thức đạo hàm giải tích trên phần cứng lượng tử thực tế mà không cần xấp xỉ sai phân hữu hạn:"),
                       ("Công thức Shift-Rule:", "∂⟨M⟩ / ∂θ_j = [⟨M⟩(θ_j + s) − ⟨M⟩(θ_j − s)] / (2 sin s) với bước dịch chuyển s = π/2."),
                       ("Bản chất vật lý:", "Parameter-Shift Rule là chuẩn vàng thực nghiệm trên máy lượng tử thật.")
                   ])

    add_clean_card(s21, Inches(6.8), Inches(1.55), Inches(5.733), Inches(5.1),
                   "Kết quả Đo lường Thực nghiệm & Kết luận",
                   [
                       ("Thiết lập kiểm thử đối chứng:", "Khởi tạo ngẫu nhiên góc θ, nạp các patch ảnh thực tế từ BreastMNIST và OCTMNIST, tính song song ma trận Jacobian bằng cả hai phương pháp."),
                       ("Sai lệch tuyệt đối đo được:", "|Δ| = |g_backprop − g_shift| < 4.1 × 10⁻⁸ trên toàn bộ các thành phần gradient."),
                       ("Phân tích sai số:", "Sai lệch hoàn toàn nằm trong giới hạn sai số làm tròn số học dấu phẩy động (float32/64)."),
                       ("Kết luận thẩm định:", "Khẳng định thuật toán lan truyền ngược mô phỏng phản ánh chính xác 100% động học vật lý lượng tử trước khi huấn luyện.")
                   ])

    # =========================================================================
    # SLIDE 22: BACKUP 4 - LIVE DEFENSE DEMO DOSSIER
    # =========================================================================
    s22 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s22, C_WHITE)
    add_header(s22, "Phụ lục Phòng vệ 4: Hệ thống Demo Trực quan",
               "Hồ sơ hệ thống chạy demo trực tiếp (gd4_defense_demo.ipynb) và video dự phòng 104 giây")
    add_footer(s22, 22, note="Slide Phụ trợ Trả lời Phản biện (Backup Slide 4)")

    add_clean_card(s22, Inches(0.8), Inches(1.55), Inches(5.7), Inches(5.1),
                   "Cấu trúc Notebook Demo (gd4_defense_demo.ipynb)",
                   [
                       ("Tác vụ thực hiện:", "Nạp trực tiếp một ảnh siêu âm vú từ test set, trích xuất feature map LIVE bằng 3 mạch lượng tử (basic_L2, strongly_L2, random_L2)."),
                       ("Huấn luyện đối xứng:", "Huấn luyện classifier head đối xứng ngay trong notebook (20 epochs, ~1.3 giây trên đặc trưng precompute)."),
                       ("Dự đoán & Chi phí:", "Xuất xác suất dự đoán ung thư vú và biểu đồ so sánh chi phí suy luận CPU giữa Classical CNN và Quanvolution."),
                       ("Mẫu thử nghiệm công khai:", "Ảnh kiểm thử được chọn theo tiêu chí khách quan in rõ trong notebook: trường hợp ác tính được phân loại đúng với độ tin cậy cao nhất (idx 6, p = 0.884)."),
                       ("Độ trung thực tính toán:", "Sai lệch giữa đặc trưng trích xuất LIVE và tensor precompute: tối đa 1.47 × 10⁻⁸.")
                   ])

    add_clean_card(s22, Inches(6.8), Inches(1.55), Inches(5.733), Inches(5.1),
                   "Video Dự phòng & Khả năng Chuyển giao",
                   [
                       ("Video minh họa (demo_defense_backup.mp4):", "Thời lượng đúng 104 giây, độ phân giải nét cao, có phụ đề tiếng Việt rõ ràng, phòng ngừa sự cố mạng hoặc gián đoạn kỹ thuật tại phòng bảo vệ."),
                       ("Script tái tạo (render_demo_video.py):", "Tự động render lại video từ đầu khi có bất kỳ thay đổi nào trong mã nguồn pipeline."),
                       ("Mã nguồn mở độc lập:", "Toàn bộ tài nguyên demo nằm trong thư mục notebooks/ của repository GitHub, có thể clone và chạy lại trên bất kỳ máy tính nào có Python 3.10."),
                       ("Tính tương thích phần mềm:", "Được kiểm thử thành công trên cả môi trường JupyterLab, VSCode Notebook và terminal console.")
                   ])

    # Save
    if not output_path:
        output_path = r'd:\KhoaLuanTotNghiep\slides\KLTN_slides.pptx'
    prs.save(output_path)
    print(f"Successfully generated presentation: {output_path}")

if __name__ == '__main__':
    out_dir = r'd:\KhoaLuanTotNghiep\slides'
    os.makedirs(out_dir, exist_ok=True)

    # Deck 1: Under Review (Default)
    path_under_review = os.path.join(out_dir, 'KLTN_slides_under_review.pptx')
    build_presentation(is_accepted=False, output_path=path_under_review)

    # Deck 2: Accepted Template (Prepared layout with reserved slot)
    path_accepted = os.path.join(out_dir, 'KLTN_slides_accepted_template.pptx')
    build_presentation(is_accepted=True, output_path=path_accepted)

    print("ALL DECKS RE-GENERATED WITH CONCISE, BALANCED LAYOUT SUCCESSFULLY!")
