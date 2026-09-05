# -*- coding: utf-8 -*-
"""generate_slides_faculty.py — WS-D v2: slide bảo vệ trên TEMPLATE KHOA
(KTMT_Phụ lục 6), nội dung cô đọng theo chuẩn academic-pptx:
1 ý/slide · ≤3 thẻ/slide · mỗi thẻ ≤35 từ · stat callout nổi bật · giữ branding CE-UIT."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[2]
TPL = ROOT / 'KTMT_KLTN_Phu luc 6_Mau bao cao bao ve.pptx'
OUT = ROOT / 'slides' / 'KLTN_slides_faculty.pptx'
NAVY = RGBColor(0x1F, 0x4E, 0x79)
ACC = RGBColor(0x2E, 0x75, 0xB6)
DARK = RGBColor(0x26, 0x26, 0x26)

base = Presentation(TPL)
SW, SH = base.slide_width, base.slide_height
blank = next(l for l in base.slide_masters[0].slide_layouts if l.name == 'Blank')

# xóa 15 slide mẫu của template, giữ master/layout/theme
for s in list(base.slides._sldIdLst):
    base.slides._sldIdLst.remove(s)

def bg_white(s):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

def txt(s, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font='Arial'):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb

def card(s, x, y, w, h, title, bullets, stat=None):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.name = 'Arial'; r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = NAVY
    for b in bullets:
        pb = tf.add_paragraph(); rb = pb.add_run(); rb.text = '• ' + b
        rb.font.name = 'Arial'; rb.font.size = Pt(13.5); rb.font.color.rgb = DARK
        pb.space_before = Pt(6)
    if stat:
        ps = tf.add_paragraph(); ps.space_before = Pt(10)
        rs = ps.add_run(); rs.text = stat
        rs.font.name = 'Arial'; rs.font.size = Pt(28); rs.font.bold = True; rs.font.color.rgb = ACC

def slide(title, section, cards, page):
    s = base.slides.add_slide(blank)
    bg_white(s)
    txt(s, 0.55, 0.28, 12.2, 0.4, section, 12, color=ACC)
    txt(s, 0.55, 0.62, 12.2, 0.75, title, 26, bold=True, color=NAVY)
    n = len(cards)
    w = 12.23 / n
    for i, cd in enumerate(cards):
        t, bullets = cd[0], cd[1]
        stat = cd[2] if len(cd) > 2 else None
        card(s, 0.55 + i * (w + 0.06), 1.55, w, 5.2, t, bullets, stat)
    txt(s, 11.9, 7.08, 0.9, 0.3, f'{page} / 22', 11, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.RIGHT)
    return s

# ============ 18 SLIDE CHÍNH (cô đọng) ============
S = []
S.append(('Bài báo SOICT 2026 (Springer CCIS) — đã nộp, đang bình duyệt', 'Phần 1: Bối cảnh',
    [('Đề tài', ['Quanvolution 4-qubit phân loại ảnh y tế MedMNIST',
                 'Benchmark đối xứng 1:1 với CNN cổ điển',
                 '10 seeds × 20 epochs × 6 metrics y tế'], None)], 1))
S.append(('Ba khoảng trống của literature QML y tế', 'Phần 1: Bối cảnh',
    [('L1 — Baseline', ['CNN tùy tiện, chênh hàng nghìn tham số', 'Không cô lập tầng lượng tử'], None),
     ('L2 — Thiếu thống kê', ['Chỉ 1–3 runs', 'Không CI / Wilcoxon / effect size'], None),
     ('L3 — Fixed vs Trainable', ['Chưa lượng hóa trên cùng thiết lập'], None)], 2))
S.append(('Nguyên lý Quanvolution: mạch lượng tử = bộ lọc ảnh', 'Phần 2: Phương pháp',
    [('Cơ chế', ['Patch 2×2 → mã hóa góc RY(πx) → 4 qubit', 'Đo ⟨Z⟩ → 4 kênh feature map [−1,1]'],
      '196 patches/ảnh')], 3))
S.append(('Ba họ ansatz — trục so sánh trung tâm', 'Phần 2: Phương pháp',
    [('Basic', ['1 trục RY + CNOT vòng', '4L tham số'], '8 (L2)'),
     ('Strongly', ['3 trục U₃ + rối mở rộng', '12L tham số'], '24 (L2)'),
     ('Random', ['Haar-random, đóng băng', '0 tham số'], '0')], 4))
S.append(('Đối xứng 1:1 — head cổ điển giống hệt nhau', 'Phần 2: Phương pháp',
    [('Head cố định', ['BN(4) → ReLU → Linear(784→K)', 'K=2: 1,570 · K=4: 3,140'],
      'Chênh kernel: 20 tham số'), ('Ý nghĩa', ['Chênh lệch hiệu năng quy hết về phép biến đổi đặc trưng'], None)], 5))
S.append(('6 metrics cho dữ liệu y tế lệch lớp', 'Phần 2: Phương pháp',
    [('Ngưỡng', ['Accuracy (tham khảo)', 'Balanced Acc', 'F1 macro', 'MCC'], None),
     ('Xếp hạng', ['ROC-AUC (OvR macro)', 'PR-AUC — metric lâm sàng số 1'], None)], 6))
S.append(('Kiến trúc phần mềm 4 tầng + kiểm chứng đạo hàm', 'Phần 2: Phương pháp',
    [('4 tầng', ['Dữ liệu → Mô hình → Thí nghiệm → Đầu ra', 'Mọi số qua kiểm định tự động'],
      '|Δ| < 4.1×10⁻⁸')], 7))
S.append(('BreastMNIST: mạch tĩnh vượt CNN (10 seeds)', 'Phần 3: Kết quả',
    [('Quán quân ROC-AUC', ['Fixed Basic: 0.8521 ± 0.0095', 'CNN: 0.8336 ± 0.0259'],
      'p = 0.0298 · d = +0.815'),
     ('Quán quân PR-AUC', ['Fixed Strongly: 0.9182 ± 0.0071', 'd = +1.332 so với CNN'], None)], 8))
S.append(('Ổn định phương sai ~2.7× — structural regularizer', 'Phần 3: Kết quả',
    [('Giải thích', ['std 0.0095 vs 0.0259 (tỷ 2.7263×)', 'Mạch tĩnh = bộ điều hòa cấu trúc, chống overfit 546 mẫu'],
      '~2.7×')], 9))
S.append(('OCTMNIST: CNN áp đảo toàn diện (10 seeds)', 'Phần 3: Kết quả',
    [('CNN dẫn đầu 6/6 metrics', ['ROC-AUC 0.7505 ± 0.0240', 'BAcc 0.4433 ± 0.0135'],
      'd = +2.108'), ('Nguyên nhân', ['Expressibility bottleneck 4-qubit L=1'], None)], 10))
S.append(('Nội bộ strongly: tự học > tĩnh, hòa quán quân', 'Phần 3: Kết quả',
    [('Cùng họ', ['Trainable 0.6922 ± 0.0199 vs Fixed 0.6690 ± 0.0055'], 'Δ+0.0232 · p=0.0098'),
     ('Vs quán quân random_L1', ['0.6912 ± 0.0071 — hòa (p=0.8875)'], None)], 11))
S.append(('Gradient khỏe 0.2–0.5 — không Barren Plateaus', 'Phần 3: Kết quả',
    [('Sanity check', ['Seed-mean 0.2–0.5, đỉnh ~1.3', 'Mạch 4-qubit nông'], '0.2–0.5')], 12))
S.append(('Chi phí: precompute là bắt buộc', 'Phần 3: Kết quả',
    [('CPU suy luận', ['CNN 0.31 ms vs Quanv 220 ms'], '~710×')], 13))
S.append(('Ba kết luận trung thực', 'Phần 4: Kết luận',
    [('Kết luận', ['Ưu thế phụ thuộc chế độ dữ liệu', 'Mạch tĩnh 0 tham số = inductive bias',
                   'Trainability chỉ cục bộ trong cùng họ'], None)], 14))
S.append(('Bài báo SOICT 2026 — đã nộp, đang bình duyệt', 'Phần 4: Đóng góp',
    [('4 đóng góp', ['Benchmark đối xứng + data-regime boundary', 'Tài liệu tái lập + bản thảo quốc tế'],
      'Springer CCIS')], 15))
S.append(('4 hướng phát triển', 'Phần 4: Kết luận',
    [('Tiếp theo', ['GPU cuQuantum', 'NISQ IBM thật', 'Full OCTMNIST 97k', 'Encoding HUAR/Hamiltonian'], None)], 16))
S.append(('Demo live: ảnh → feature map → dự đoán', 'Phần 5: Demo',
    [('Live', ['Train head 1.3s trên feature precompute', 'Malignant idx 6 — p = 0.884'],
      '1.47×10⁻⁸')], 17))
S.append(('Liêm chính & tái lập', 'Phần 5: Cam kết',
    [('Bằng chứng', ['Mọi số truy vết JSON gốc', 'Tag soict-submission-v4',
                     '19/19 refs đã xác minh'], '92/92 ô')], 18))

# dựng 18 slide chính
for i, (title, section, cards, page) in enumerate(S, 1):
    slide(title, section, cards, page)

# Q&A + 4 backup slides (giữ nguyên vai trò như deck cũ)
s = base.slides.add_slide(blank); bg_white(s)
txt(s, 0.5, 2.8, 12.33, 1.2, 'Q&A — Trân trọng cảm ơn Hội đồng', 36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txt(s, 0.5, 4.2, 12.33, 0.6, 'Mã nguồn & dữ liệu: github.com/NamIsStudyingCE/Quanvolution (tag soict-submission-v4)',
    14, color=ACC, align=PP_ALIGN.CENTER)

BACKUP = [
    ('Backup 1 — Kiểm định thống kê đầy đủ', [
        ('Breast CNN vs Fixed Basic (AUC)', ['Δ = −0.0186 · p_t = 0.0298 · p_w = 0.0254 · d = −0.815']),
        ('Breast CNN vs Fixed Strongly (PR-AUC)', ['Δ = −0.0140 · p_t = 0.0023 · p_w = 0.0059 · d = −1.332']),
        ('OCT CNN vs Trainable Strongly (AUC)', ['Δ = +0.0583 · p_t ≈ 0.0001 · p_w = 0.0020 · d = +2.108']),
        ('OCT Trainable vs Fixed Strongly (AUC)', ['Δ = +0.0232 · p_t = 0.0090 · p_w = 0.0098 · d = +1.050'])]),
    ('Backup 2 — Phân rã tham số đối xứng 1:1', [
        ('BreastMNIST K=2', ['CNN: 20 kernel + 8 BN + 1,570 head = 1,598', 'Fixed: 0 kernel + 8 BN + 1,570 head = 1,578']),
        ('OCTMNIST K=4', ['CNN: 3,168 · Fixed: 3,148 · Trainable Strongly: 3,160'])]),
    ('Backup 3 — Kiểm chứng đạo hàm Parameter-Shift', [
        ('Hai đường tính gradient', ['Analytic backprop (default.qubit)', 'Parameter-Shift Rule: ∂F/∂θ = [F(θ+π/2) − F(θ−π/2)]/2']),
        ('Kết quả', ['|Δ| < 4.1×10⁻⁸ — sai số số học float'])]),
    ('Backup 4 — Hồ sơ demo dự phòng', [
        ('Demo notebook', ['Ảnh test idx 6 (malignant) → p = 0.884 — phân loại đúng',
                           'Live vs precompute: sai lệch 1.47×10⁻⁸']),
        ('Video dự phòng', ['demo_defense_backup.mp4 — 104 giây, phụ đề tiếng Việt'])]),
]
for title, cards in BACKUP:
    slide(title, 'Phụ lục phòng vệ', cards, 22)

base.save(OUT)
print('saved', OUT, '| slides:', len(base.slides.__iter__.__self__._sldIdLst))
