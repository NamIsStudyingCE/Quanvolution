# -*- coding: utf-8 -*-
"""fill_faculty_template.py — WS-D v3: điền nội dung khóa luận VÀO TEMPLATE KHOA
bằng cách thay text in-place trong placeholder (giữ 100% thiết kế/layout).
v3: nội dung đơn giản hóa (ít từ chuyên ngành, câu ngắn) + 4 slide trống
chỉ còn chrome (badge/logo + vòng tròn số trang góc trái + Copyright) để
chủ dự án tự chèn bảng bằng tay."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
from copy import deepcopy

TPL = Path(__file__).resolve().parents[2] / 'KTMT_KLTN_Phu luc 6_Mau bao cao bao ve.pptx'
OUT = Path(__file__).resolve().parents[2] / 'slides' / 'KLTN_slides_faculty_inplace.pptx'

prs = Presentation(str(TPL))
slides = list(prs.slides)

def set_tf(tf, lines, size=None, bold=None):
    tf.clear()
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln
        if size: r.font.size = Pt(size)
        if bold is not None: r.font.bold = bold

def shape_by_text(slide, starts):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip().startswith(starts):
            return sh
    raise KeyError(f'shape "{starts}" not found')

# ---------- SLIDE 1: Bìa ----------
s1 = slides[0]
for sh in s1.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() == 'TÊN KHÓA LUẬN':
        set_tf(sh.text_frame, [
            'QUANVOLUTION — MẠCH LƯỢNG TỬ PHÂN LOẠI ẢNH Y TẾ',
            'So sánh công bằng với CNN trên MedMNIST',
        ], size=30, bold=True)
    if sh.has_text_frame and sh.text_frame.text.strip().startswith('GVHD'):
        set_tf(sh.text_frame, [
            'GVHD: TS. Nguyễn Duy Xuân Bách — Khoa Kỹ thuật Máy tính, UIT',
            'Sinh viên thực hiện: Nguyễn Hạo Nam',
            'Bài báo khoa học: đã nộp tại SOICT 2026 (Springer CCIS) — đang bình duyệt',
        ])

# ---------- SLIDE 4: Tổng quan (đơn giản hóa) ----------
s4 = slides[3]
set_tf(shape_by_text(s4, 'Giới thiệu tổng quan').text_frame, [
    'Ảnh y tế: dữ liệu nhỏ, ít ca bệnh, lệch lớp',
    'CNN lớn → dễ học thuộc dữ liệu nhỏ (overfitting)',
    'Máy học lượng tử: hướng mới, mạch 4-qubit làm "bộ lọc ảnh"',
    '3 vấn đề cũ: so sánh không công bằng · ít lần chạy · chưa rõ mạch tĩnh hay tự học',
    'Đề tài: thử nghiệm công bằng — 10 lần chạy × 6 chỉ số × 2 bộ dữ liệu',
])

# ---------- SLIDE 6: Giải pháp (lý thuyết, đơn giản hóa) ----------
s6 = slides[5]
set_tf(shape_by_text(s6, 'Phân tích lý thuyết').text_frame, [
    'Quanvolution: từng ô 2×2 của ảnh → mạch 4-qubit → đo ra 4 giá trị đặc trưng',
    'Ba loại mạch: Basic · Strongly · Random (không cần học)',
    'Khác QCNN: chỉ phần trích xuất đặc trưng là lượng tử, phần phân loại là cổ điển',
    'Đạo hàm kiểm chứng bằng 2 cách — sai lệch chỉ 0.00000004',
])

# ---------- SLIDE 7: Giải pháp (thiết kế) ----------
s7 = slides[6]
set_tf(shape_by_text(s7, 'Phân tích thiết kế').text_frame, [
    'Cùng một head phân loại cho mọi mô hình — so sánh công bằng',
    'Mạch tĩnh: tính sẵn đặc trưng 1 lần — 10 lần chạy chỉ mất ~18 giây',
    'Mạch tự học: huấn luyện xuyên qua mạch lượng tử — dùng 2 tốc độ học',
    'Quy trình 4 tầng + kiểm tra số liệu tự động',
    '6 chỉ số: Accuracy · BAcc · F1 · MCC · ROC-AUC · PR-AUC',
])

# ---------- SLIDE 9: Kết quả BreastMNIST ----------
s9 = slides[8]
set_tf(shape_by_text(s9, 'Kết quả mô phỏng').text_frame, [
    'Mạch Basic: ROC-AUC 0.8521 ± 0.0095 — cao nhất, hơn CNN 0.8336 ± 0.0259 (p = 0.0298)',
    'Mạch Strongly: PR-AUC 0.9182 ± 0.0071 — cao nhất (p = 0.0023)',
    'Ổn định hơn CNN khoảng 2.7 lần (0.0095 so với 0.0259)',
    'Mạch tự học: Balanced Acc 0.6945 ± 0.0451 — chưa khác biệt CNN (p = 0.6701)',
])

# ---------- SLIDE 10: Kết quả OCTMNIST ----------
s10 = slides[9]
set_tf(shape_by_text(s10, 'Kết quả hiển thị').text_frame, [
    'CNN dẫn đầu cả 6 chỉ số: ROC-AUC 0.7505 ± 0.0240, PR-AUC 0.4991 ± 0.0297',
    'Cách biệt với mạch lượng tử tốt nhất: Δ = +0.0583 (p ≈ 0.0001, d = +2.108)',
    'Nguyên nhân: mạch 4-qubit quá nhỏ, head chỉ có 3,140 tham số (K = 4)',
    'Chi phí: mạch lượng tử 220 ms vs CNN 0.31 ms mỗi ảnh (~710×)',
    'Demo: ảnh siêu âm → đặc trưng → dự đoán malignant p = 0.884',
])

# ---------- SLIDE 12: Kết luận & hướng phát triển ----------
s12 = slides[11]
set_tf(shape_by_text(s12, 'Kết luận').text_frame, [
    'Kết luận 1 — Cơ sở dữ liệu quyết định hiệu quả: thắng ở dữ liệu nhỏ, thua ở dữ liệu lớn',
    'Kết luận 2 — Mạch tĩnh 0 tham số vẫn cho kết quả tốt và rất ổn định',
    'Kết luận 3 — Tự học chỉ giúp trong cùng loại mạch, không hơn mạch tĩnh tối ưu',
    'Hướng phát triển: GPU · máy lượng tử thật IBM · đủ 97,000 ảnh OCT · mã hóa mới',
])

# ---------- SLIDE 13: Tài liệu tham khảo (rút gọn) ----------
s13 = slides[12]
set_tf(shape_by_text(s13, 'Tài liệu tham khảo').text_frame, [
    'Henderson et al. (2020). Quanvolutional neural networks. QMI, 2(1), 2.',
    'Cong et al. (2019). Quantum convolutional neural networks. Nature Physics, 15(12).',
    'Matondo-Mvula & Elleithy (2024). Breast cancer detection with QNN. Entropy, 26(8), 630.',
    'Schuld & Killoran (2022). Is quantum advantage the right goal for QML? PRX Quantum, 3(3), 030101.',
    'Yang et al. (2023). MedMNIST v2 benchmark. Scientific Data, 10(1), 41.',
    'Azevedo et al. (2022). Quantum transfer learning for breast cancer. QMI, 4(1), 5.',
    '… danh mục đầy đủ 19 tài liệu — xem Luận văn, Chương 2.',
])

# ---------- SLIDE 15: Phụ lục (tài liệu phòng vệ) ----------
s15 = slides[14]
set_tf(shape_by_text(s15, 'Phần phụ lục này đặt').text_frame, [
    'Tài liệu phòng vệ 1: Bảng kiểm định thống kê đầy đủ — Δ, p (t-test & Wilcoxon), Cohen\u2019s d (Bảng 4.3 luận văn)',
    'Tài liệu phòng vệ 2: Phân rã tham số K=2 (1,570) / K=4 (3,140) — Bảng 3.2 luận văn',
    'Tài liệu phòng vệ 3: Kiểm chứng đạo hàm Parameter-Shift — |Δ| < 4.1×10⁻⁸',
    'Tài liệu phòng vệ 4: Demo notebook + video dự phòng 104 giây (Phụ lục D luận văn)',
    'Mã nguồn & số liệu: github.com/NamIsStudyingCE/Quanvolution — tag soict-submission-v4',
])

# ---------- 4 SLIDE TRỐNG CHỈ CÒN CHROME (để tự chèn bảng bằng tay) ----------
donor = slides[14]  # slide Phụ lục (layout 1_Title and Content, có đủ chrome)
CHROME_N = 4
for k in range(CHROME_N):
    s = prs.slides.add_slide(donor.slide_layout)
    # xóa placeholder tiêu đề/nội dung mẫu inherited từ layout (để slide trống hoàn toàn)
    for ph in list(s.placeholders):
        if ph.placeholder_format.idx in (0, 1):
            ph._element.getparent().remove(ph._element)
    # copy chrome từ slide donor: Copyright + số trang (placeholder bên phải-dưới)
    for sh in donor.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text
        is_page = sh.left is not None and sh.left > Inches(10) and sh.top > Inches(6)  # ‹#› góc phải-dưới
        if 'Copyrights' in t or is_page:
            el = deepcopy(sh._element)
            new_sh = s.shapes._spTree.append(el)
    # chuyển số trang (vòng tròn ‹#›) xuống GÓC DƯỚI BÊN TRÁI theo yêu cầu
    for sh in s.shapes:
        if sh.has_text_frame and sh.left is not None and sh.left > Inches(10) and sh.top > Inches(6) \
           and 'Copyrights' not in sh.text_frame.text:
            sh.left = Inches(0.55)
            sh.top = Inches(6.85)

prs.save(str(OUT))
print('saved', OUT, '| tổng slide:', len(prs.slides._sldIdLst),
      f'(15 template + {CHROME_N} slide trống chrome)')
